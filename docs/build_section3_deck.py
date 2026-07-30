"""Build the Software System (requirements + architecture) crit section.

Palette and grid come from eParts_Section4_10min.pptx so the two sections read as one
deck. Everything on a slide is >= 20 pt so it survives projection; that is the binding
constraint on how much text each card can hold, so the slides carry phrases and the
talking script carries the sentences.
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

BLACK = RGBColor(0x00, 0x00, 0x00)
BANNER = RGBColor(0x11, 0x11, 0x11)
CARD = RGBColor(0x1C, 0x1C, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x99, 0x99, 0x99)      # lifted from 777777 — 20 pt grey needs contrast
LEAD = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT = WHITE                           # emphasis
DIM = RGBColor(0xAA, 0xAA, 0xAA)         # secondary emphasis, one step under ACCENT
RULE = RGBColor(0x44, 0x44, 0x44)        # hairlines

MIN_PT = 20.0                            # nothing on a slide may be smaller
FONT = "Aptos"
BASE = "/Users/arjun/Documents/CMU/studio-project/"
OUT_SOFTWARE = BASE + "eParts_Section3_SoftwareSystem.pptx"
OUT_REFLECTION = BASE + "eParts_Reflection_Closing.pptx"
DIAG = "/Users/arjun/Documents/CMU/studio-project/Diagrams/"
CONFLUENCE = ("https://cmu-mse.atlassian.net/wiki/spaces/AISDLC/pages/76742657/"
              "Engineering+System+Artifacts")
ADR_INDEX = "https://github.com/AshrithaG/eparts/blob/main/docs/adr-index.md"
# ADR-018 is the one we open live — it is the exemplar for the format. Points at
# Confluence rather than GitHub now that the ADRs are published there, so the room
# lands on a rendered page instead of raw Markdown.
ADR_018 = ("https://cmu-mse.atlassian.net/wiki/spaces/AISDLC/pages/78708738/"
           "ADR-018+Extend+Routing+to+ETIM+Signals+with+a+Class-Review-First+Path")
ADR_018_SHOWN = "cmu-mse.atlassian.net/wiki/spaces/AISDLC/pages/78708738"

def new_deck():
    """A fresh 720x405 pt presentation. Two decks are built from this one script so the
    palette, grid and 20 pt floor cannot drift between them."""
    global prs, BLANK
    prs = Presentation()
    prs.slide_width = Emu(9144000)       # 720 pt
    prs.slide_height = Emu(5143500)      # 405 pt
    BLANK = prs.slide_layouts[6]
    return prs


prs = new_deck()


def P(v):
    return Emu(int(v * 12700))


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, P(x), P(y), P(w), P(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def obox(s, x, y, w, h, color=RULE, dashed=False, fill=CARD):
    """Outlined box. The detail slide needs dashed strokes for designed-not-built."""
    sh = s.shapes.add_shape(1, P(x), P(y), P(w), P(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = color
    sh.line.width = Pt(1.25)
    if dashed:
        sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, size=MIN_PT, color=MUTED, bold=False, space=0,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.18, gap=7, wrap=True):
    assert size >= MIN_PT, f"{size} pt is below the {MIN_PT} pt projection floor"
    tb = s.shapes.add_textbox(P(x), P(y), P(w), P(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        if i:
            p.space_before = Pt(gap)
        for txt, b, c in (para if isinstance(para, list) else [(para, bold, color)]):
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = b
            f.color.rgb = c
            if space:
                r.font._rPr.set("spc", str(int(space * 100)))
    return tb


def title(s, txt):
    text(s, 28.8, 21.6, 662.4, 40, txt, size=28, color=WHITE, bold=True, line=1.0)


def artifact_link(s, label, url=None):
    tb = s.shapes.add_textbox(P(340), P(3.0), P(351.2), P(16.0))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    r = para.add_run()
    r.text = label
    r.hyperlink.address = url or CONFLUENCE
    f = r.font
    f.name, f.size, f.bold, f.underline = FONT, Pt(MIN_PT), True, False
    f.color.rgb = ACCENT
    return tb


def banner(s, chunks, y=68.0, h=52.0):
    rect(s, 28.8, y, 662.4, h, BANNER)
    text(s, 42.0, y + 5, 636.0, h - 10, [chunks], color=LEAD,
         anchor=MSO_ANCHOR.MIDDLE, line=1.14)


def card(s, x, y, w, h, bar, head, body, head_color=WHITE, head_size=22):
    rect(s, x, y, w, h, CARD)
    rect(s, x, y, w, 5.04, bar)
    pad = 12.0
    text(s, x + pad, y + 12, w - 2 * pad, 30, head, size=head_size, color=head_color,
         bold=True, line=1.0)
    text(s, x + pad, y + 12 + head_size * 1.55, w - 2 * pad, h - 34 - head_size * 1.55,
         body, gap=9)


def footer(s, txt, color=ACCENT, y=372.0, link_label=None, url=None, shown_url=None):
    """Evidence strip. With link_label/url it grows a second line carrying the artifact
    link, which is where the audience expects it — bottom left, not the header."""
    h = 24.0 if not link_label else 50.0
    if link_label and y == 372.0:
        y = 346.0
    rect(s, 28.8, y, 4.5, h, color)
    text(s, 42.0, y, 649.2, 24.0, txt.upper(), color=color, bold=True, space=0.5,
         anchor=MSO_ANCHOR.MIDDLE, line=1.0)
    if not link_label:
        return
    tb = s.shapes.add_textbox(P(42.0), P(y + 25.0), P(649.2), P(24.0))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.line_spacing = 1.0
    lab = para.add_run()
    lab.text = link_label + "  "
    lab.font.name, lab.font.size, lab.font.bold = FONT, Pt(MIN_PT), True
    lab.font.color.rgb = WHITE
    ln = para.add_run()
    ln.text = shown_url or url
    ln.hyperlink.address = url
    ln.font.name, ln.font.size, ln.font.bold = FONT, Pt(MIN_PT), False
    ln.font.color.rgb = MUTED
    ln.font.underline = False


def save(deck, path, expected_slides):
    """Write the deck, then repaint theme hyperlink colours.

    python-pptx cannot set them, and the default Office theme forces links to blue,
    which is unreadable on this black background.
    """
    import re as _re, shutil as _shutil, zipfile as _zipfile
    assert len(deck.slides._sldIdLst) == expected_slides, \
        f"{path}: expected {expected_slides} slides, got {len(deck.slides._sldIdLst)}"
    deck.save(path)
    tmp = path + ".tmp"
    with _zipfile.ZipFile(path) as zin, \
         _zipfile.ZipFile(tmp, "w", _zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename.startswith("ppt/theme/theme"):
                t = data.decode("utf-8")
                # Links white, and the stock Office accent hues greyed out. Nothing on a
                # slide uses them, but a themed shape inserted later would pick up orange.
                greys = {"hlink": "FFFFFF", "folHlink": "FFFFFF", "dk2": "1C1C1C",
                         "lt2": "EEEEEE", "accent1": "FFFFFF", "accent2": "CCCCCC",
                         "accent3": "AAAAAA", "accent4": "888888", "accent5": "666666",
                         "accent6": "444444"}
                for tag, val in greys.items():
                    t = _re.sub(rf'(<a:{tag}>).*?(</a:{tag}>)',
                                rf'\1<a:srgbClr val="{val}"/>\2', t, flags=_re.S)
                data = t.encode("utf-8")
            zout.writestr(info, data)
    _shutil.move(tmp, path)
    print(f"wrote {path} — {expected_slides} slide(s), all text >= {MIN_PT} pt")


# two-column geometry
LX, RX, CW = 28.8, 368.4, 322.8
CY, CH = 126.0, 238.0

# ─────────────────────────────────────────────────────────── 1. requirements v1.0 → v1.4
# Counts verified against product-spec-v1.4.tex. v1.0 baseline: HLR-1..5, FR-1..8,
# DR-1..3, C-1..3 = 19. v1.4 adds HLR-6, FR-9, FR-10, DR-4 (v1.1) and C-4 (v1.2) = 24.
# 2 rewritten (HLR-2, FR-1). QAS-3/VAL-4/VAL-5 also arrived in v1.4 but are deliberately
# left off this table — Arjun narrates them instead of adding two more rows.
s = new_slide()
title(s, "Requirements: v1.0 → v1.4")
# ── comparison table (left)
TX, TW, TY = 28.8, 392.0, 86.0
rect(s, TX, TY, TW, 232.0, CARD)
COL_L, COL_A, COL_B = TX + 16, TX + 236, TX + 314
text(s, COL_L, TY + 10, 200, 26, "Requirement set", color=MUTED, bold=True)
text(s, COL_A, TY + 10, 70, 26, "v1.0", color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
text(s, COL_B, TY + 10, 62, 26, "v1.4", color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
rect(s, COL_L, TY + 38, TW - 32, 1.2, RULE)

ROWS = [("High-level", "5", "6"), ("Functional", "8", "10"),
        ("Derived", "3", "4"), ("Constraints", "3", "4")]
yy = TY + 46
for label, a, b in ROWS:
    text(s, COL_L, yy, 200, 26, label, color=WHITE)
    text(s, COL_A, yy, 70, 26, a, color=MUTED, align=PP_ALIGN.RIGHT)
    text(s, COL_B, yy, 62, 26, b, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    yy += 32
rect(s, COL_L, yy + 2, TW - 32, 1.2, RULE)
text(s, COL_L, yy + 10, 200, 26, "Total", color=WHITE, bold=True)
text(s, COL_A, yy + 10, 70, 26, "19", color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
text(s, COL_B, yy + 10, 62, 26, "24", color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

# ── the two numbers that matter (right)
def stat(x, y, w, h, big, caption, big_colour):
    rect(s, x, y, w, h, CARD)
    rect(s, x, y, w, 5.04, big_colour)
    text(s, x + 16, y + 12, w - 32, 48, big, size=40, color=big_colour, bold=True,
         line=1.0)
    text(s, x + 16, y + 62, w - 32, 30, caption, color=MUTED)

SX, SW = 440.0, 251.2
stat(SX, TY, SW, 112.0, "89%", "of v1.0 unchanged", WHITE)
stat(SX, TY + 120.0, SW, 112.0, "37%", "churn since April", DIM)
footer(s, "5 new IDs · 2 rewritten · 17 of 19 untouched",
       link_label="Spec v1.4:", url=CONFLUENCE, shown_url="cmu-mse.atlassian.net/wiki/spaces/AISDLC")

# ─────────────────────────────────────────── 2. quality attributes
# Sits next to the requirements slide because that is what these are. Attribute name in
# the middle column, the measure in plain words on the right — a QAS whose measure you
# can't say out loud isn't doing any work. Two of the four are modifiability, which is
# the point: it was the April bet and ETIM is what tested it.
s = new_slide()
title(s, "Quality attributes driving the design")

QID_X, QATT_X, QM_X = 42.0, 116.0, 300.0
HDR_Y = 84.0
text(s, QID_X, HDR_Y, 68, 26, "QAS", color=MUTED, bold=True)
text(s, QATT_X, HDR_Y, 170, 26, "Attribute", color=WHITE, bold=True)
text(s, QM_X, HDR_Y, 388, 26, "What it has to do", color=MUTED, bold=True)
rect(s, 28.8, HDR_Y + 30, 662.4, 1.2, RULE)

# Numeric order. Listing them by priority put QAS-4 above QAS-3, which reads as a
# mistake rather than as emphasis — the slide is an artifact index, so it sorts by ID.
qas = [
    ("QAS-1", "Modifiability", "A new supplier format in 4 hours"),
    ("QAS-2", "Usability", "10 review decisions a minute"),
    ("QAS-3", "Modifiability", "Policy change with no code deploy"),
    ("QAS-4", "Accuracy", "95% of auto-accepted values correct"),
]
qy = HDR_Y + 42
for qid, attr, measure in qas:
    rect(s, 28.8, qy, 662.4, 44.0, CARD)
    rect(s, 28.8, qy, 4.5, 44.0, WHITE)
    text(s, QID_X, qy, 68, 44.0, qid, color=DIM, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, line=1.0, wrap=False)
    text(s, QATT_X, qy, 170, 44.0, attr, color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, line=1.0, wrap=False)
    text(s, QM_X, qy, 388, 44.0, measure, color=MUTED,
         anchor=MSO_ANCHOR.MIDDLE, line=1.0, wrap=False)
    qy += 52.0
footer(s, "QAS-4 restored from April · not yet measured",
       color=DIM, link_label="Spec v1.4:", url=CONFLUENCE,
       shown_url="cmu-mse.atlassian.net/wiki/spaces/AISDLC")

# ─────────────────────────────────────────────────────────── 3. managing the change
s = new_slide()
title(s, "How we managed the change")
banner(s, [
    ("Spec 1.0 → 1.1 → 1.2 → 1.3 → 1.4.", True, WHITE),
    ("  Each version-history entry is the change record.", False, LEAD),
])
card(s, LX, CY, 662.4, CH, WHITE, "What we did", [
    [("Added new IDs instead of editing old ones — HLR-6, FR-9/10, DR-4, C-4, QAS-3.",
      False, MUTED)],
    [("Renumbered nothing, so every trace link from April still resolves.", False, MUTED)],
    [("Every ETIM decision traces to a requirement, and forward to code and a test.",
      False, MUTED)],
])
footer(s, "HLR-6 → FR-9 → ADR-16 → ETIM tables in code → 10 tests")

# ─────────────────────────────────────────────────────────── 3. architecture v5 → v6
s = new_slide()
title(s, "How the architecture changed")
IY = 84.0
# v6 is v5 plus one box, so the two thumbnails are near-identical in aspect (1.09 and
# 1.01). Same height, same baseline — the eye is meant to compare silhouettes.
s.shapes.add_picture(DIAG + "pipe-filter-architecturev5-grey.png", P(28.8), P(IY + 24),
                     width=P(163), height=P(150))
s.shapes.add_picture(DIAG + "pipe-filter-architecture-v6-grey.png", P(199), P(IY + 24),
                     width=P(151), height=P(150))
text(s, 28.8, IY + 182, 163, 24, "v5.0 · MAY", color=MUTED, bold=True)
text(s, 199, IY + 182, 151, 24, "v6.0 · JULY", color=ACCENT, bold=True)

DX, DW = 362.0, 329.2
rect(s, DX, IY, DW, 228.0, CARD)
rect(s, DX, IY, DW, 5.04, ACCENT)
text(s, DX + 12, IY + 12, DW - 24, 28, "What ETIM changed", size=22, color=ACCENT,
     bold=True)
yy = IY + 48
for n, label in [("1", "ETIM matching, a new phase"),
                 ("2", "ETIM dictionary loaded"),
                 ("3", "Raw values kept separately"),
                 ("4", "Fixed handoff format to ML"),
                 ("5", "PIMS keyed by ETIM IDs")]:
    text(s, DX + 12, yy, 16, 26, n, color=ACCENT, bold=True)
    text(s, DX + 30, yy, DW - 44, 26, label, color=WHITE)
    yy += 33
footer(s, "one new box · nothing else moved",
       link_label="Full v6.0 diagram:", url=CONFLUENCE, shown_url="cmu-mse.atlassian.net/wiki/spaces/AISDLC")

# ─────────────────────────────────────────── 4. the one change that needs explaining
# The v6 diagram is portrait 1600x1898. At full slide height it is 253 pt wide and its
# stage labels land near 1.4 pt, so neither scaling nor cropping the bitmap can make it
# readable on a projector. This slide redraws the changed region as native shapes so
# every label is 20 pt. Slide 3 keeps the thumbnails, whose only job is silhouette
# comparison — you don't have to read those to see that the shape didn't move.
s = new_slide()
title(s, "Inside the new box: two passes")

BW, BG = 121.0, 14.35           # five boxes across 662.4 pt of usable width
BX0, ROW2_Y, BH = 28.8, 236.0, 62.0

# ── pass one: what already existed
obox(s, BX0, 108.0, 662.4, 52.0, color=WHITE, fill=CARD)
text(s, BX0 + 16, 108.0, 400, 52.0, "ML attribute matching", color=WHITE, bold=True,
     anchor=MSO_ANCHOR.MIDDLE, line=1.0)
text(s, BX0 + 430, 108.0, 216, 52.0, "already existed", color=MUTED,
     anchor=MSO_ANCHOR.MIDDLE, line=1.0, align=PP_ALIGN.RIGHT)

# ── the seam
text(s, BX0, 168.0, 662.4, 26, "then, in the same service", color=DIM, line=1.0)
text(s, BX0, 200.0, 400, 26, "ML ETIM matching  ·  NEW", color=WHITE, bold=True,
     space=0.6, line=1.0)

# ── pass two: five stages, dashed because none of this is built yet
for i, label in enumerate(["Class", "Feature", "Value\n+ unit", "ETIM\ncheck",
                           "Policy\ncheck"]):
    bx = BX0 + i * (BW + BG)
    obox(s, bx, ROW2_Y, BW, BH, color=DIM, dashed=True, fill=None)
    text(s, bx + 6, ROW2_Y, BW - 12, BH, label.split("\n"), color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line=1.0, gap=0)
    if i < 4:
        text(s, bx + BW, ROW2_Y, BG, BH, "→", color=DIM,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line=1.0)

text(s, BX0, ROW2_Y + BH + 14, 662.4, 26,
     [[("Wrong class, wrong feature list.", True, WHITE),
       (" Every attribute under it is wrong.", False, MUTED)]],
     line=1.0)
footer(s, "dashed = designed, not built", color=DIM,
       link_label="ADR-018:", url=ADR_018,
       shown_url="cmu-mse.atlassian.net/wiki/spaces/AISDLC")

# ─────────────────────────────────────────────────────────── 5. decisions
# The point of this slide is the pairing: each decision beside the alternative we
# turned down. A reason column would only restate the decision.
#
# The ADR column is what turns four claims into four artifact references, which is
# what the rubric's "architectural descriptions and decisions" is asking for. Row 3
# has no single ADR — the artifacts *are* 16 through 21 — so it cites the range
# rather than inventing a number.
s = new_slide()
title(s, "Decisions, and what we chose against")

ADR_X, CHOSE_X, VS_X = 42.0, 116.0, 440.0
HDR_Y = 84.0
text(s, ADR_X, HDR_Y, 68, 26, "ADR", color=MUTED, bold=True)
text(s, CHOSE_X, HDR_Y, 296, 26, "We chose", color=WHITE, bold=True)
text(s, VS_X, HDR_Y, 248, 26, "Instead of", color=MUTED, bold=True)
rect(s, 28.8, HDR_Y + 30, 662.4, 1.2, RULE)

decisions = [
    ("16", "ML does the matching", "ETIM keys at normalization"),
    ("14", "Raw values in own table", "one table holding both"),
    ("16–21", "New ADRs, not edits", "editing the April ones"),
    ("20", "Stay on ETIM 10.0", "building an upgrade path"),
]
ry, RH, RG = HDR_Y + 42, 44.0, 8.0
for adr, chose, instead in decisions:
    rect(s, 28.8, ry, 662.4, RH, CARD)
    rect(s, 28.8, ry, 4.5, RH, WHITE)
    # wrap=False on both text cells: a wrapped ADR range dragged the decision onto
    # two lines, and neither column is ever long enough to need wrapping.
    text(s, ADR_X, ry, 68, RH, adr, color=DIM, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, line=1.0, wrap=False)
    text(s, CHOSE_X, ry, 296, RH, chose, color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, line=1.0, wrap=False)
    text(s, VS_X - 24, ry, 18, RH, "×", color=RGBColor(0x66, 0x66, 0x66),
         anchor=MSO_ANCHOR.MIDDLE, line=1.0)
    text(s, VS_X, ry, 248, RH, instead, color=MUTED, anchor=MSO_ANCHOR.MIDDLE,
         line=1.0, wrap=False)
    ry += RH + RG
footer(s, "matching stages designed, not yet built", color=DIM,
       link_label="ADR-018:", url=ADR_018,
       shown_url="cmu-mse.atlassian.net/wiki/spaces/AISDLC")

NOTES_SOFTWARE = [
    "ETIM is a mid-project requirements CHANGE, not a new project.\n\n"
    "One line: we went from predicting free-form attributes to classifying each product "
    "into an ETIM class and matching its attributes to a controlled vocabulary.\n\n"
    "Principle: original supplier data is evidence, ETIM is a standardized "
    "interpretation over it, confidence is how sure we are of the interpretation.\n\n"
    "C-4 is the newest piece: we are pinned to ETIM 10.0 EI for the project. Adopting "
    "later releases is out of scope. If asked why — the upgrade path is a diff report, "
    "a bulk re-match and a second review queue for an event that won't happen inside "
    "this project, and nobody has decided who authorizes an upgrade. A stated limit is "
    "defensible; a half-built upgrade path isn't.\n\n"
    "Artifacts: docs/product-spec-v1.2.pdf, product-spec-changelog.md, "
    "etim-requirements-change.md.",

    "Four classes of requirements management: change control, version control, status "
    "tracking, tracing.\n\n"
    "Strongest point is version control: we ADDED IDs (HLR-6, FR-9, FR-10, DR-4, then "
    "C-4) rather than renumbering, so every existing trace link survives.\n\n"
    "We used the same discipline twice — v1.1 integrated ETIM, v1.2 pinned the release "
    "rather than silently editing v1.1.\n\n"
    "Be honest about the blocked items. ETIMARTCLASSFEATUREMAP.csv genuinely has no "
    "'required' column, so until the client defines a feature policy, 'what blocks "
    "publish?' is unanswerable. ADR-019 records the seam so the build doesn't stall.\n\n"
    "Known defect to own before it's found: two spec lineages exist — this one "
    "(0.1 -> 0.5 -> 1.0 -> 1.1 -> 1.2) and a 'Document Version 2.0' from April with a "
    "different ID set. ADRs 0001-0015 cite the April IDs; 0016-0021 cite v1.2. "
    "Recorded in product-spec-changelog.md and section 10.10 of the matrix.",

    "Open the full v6 PNG rather than squinting at the thumbnail:\n"
    "Diagrams/pipe-filter-architecture-v6.png\n\n"
    "Lead with what did NOT change — pipe-and-filter, per-attribute routing, "
    "human-in-the-loop, audit trail. The spine survived a major requirements change.\n\n"
    "Delta 2 is the one to dwell on. You can't match an attribute until you know the "
    "class, because the legal feature set is defined per class. Get the class wrong and "
    "every feature under it is wrong at high confidence, so routing won't catch it. "
    "Hence five stages and a class-review step ahead of attribute routing.\n\n"
    "Delta 3: supplier text is evidence, ETIM is interpretation. One table mixes them "
    "and you lose the ability to trace a published value back to its source.\n\n"
    "Built and merged: reference layer (alembic 0005), evidence staging split (0006), "
    "extracted_inputs handoff (0007). Not built: matching stages, ETIM-aware routing, "
    "re-keyed writeback. Say so.\n\n"
    "Telemetry question: Datadog is the production target (ADR-012 stands). Prometheus "
    "+ OpenTelemetry + structlog is the local substrate. The assessment doc read the "
    "code without the deployment intent and called it a contradiction.",

    "Four decisions, each with the alternative we rejected.\n\n"
    "Row 3 is the one I'd defend hardest: we did not edit ADRs 0001-0012. They record "
    "what we decided in April. We superseded forward and wrote an assessment that goes "
    "ADR by ADR through what ETIM affected. Editing in place erases the history.\n\n"
    "Row 4 is the newest decision. ETIM 10.0 is pinned via C-4. We accept the catalog "
    "goes stale relative to ETIM; that's the trade. The etim_release_id columns stay in "
    "the schema for provenance, so a published row still names the release it was "
    "matched under, and un-pinning later is a scope change rather than a migration.\n\n"
    "The five still-open client decisions: phase-one class list, feature policy per "
    "class, 'ETIM Other' handling, metric-canonical storage and display units, and "
    "mapping sign-off ownership. Release-upgrade governance used to be a sixth; C-4 "
    "closed it.\n\n"
    "If asked what we'd do differently: reconcile the two spec lineages earlier, and "
    "wire the handoff builder into the orchestrator (EPARTS-363) so the boundary is "
    "exercised in production flow, not only in unit tests.",
]


# ─────────────────────────────────────────────────────────── notes, deck 1
for slide, note in zip(prs.slides, NOTES_SOFTWARE):
    slide.notes_slide.notes_text_frame.text = note
save(prs, OUT_SOFTWARE, 6)


# ══════════════════════════════════════════════════════════════════════════════
# DECK 2 — Reflection & Closing. Its own file because it is presented in the
# closing slot at the end of the team talk, not as part of Software System.
# ══════════════════════════════════════════════════════════════════════════════
prs = new_deck()
s = new_slide()
title(s, "Two lessons")
LY, LH = 84.0, 262.0
card(s, LX, LY, CW, LH, DIM, "3-day cycles didn't hold", [
    [("Too much changed per tick", False, MUTED)],
    [("A small blocker ate the whole tick", False, MUTED)],
    [("Reviewing an agent PR took longer than writing it", False, MUTED)],
    [("Now ", False, MUTED), ("7-day cycles", True, WHITE), (" on Jira", False, MUTED)],
], head_color=DIM)
card(s, RX, LY, CW, LH, ACCENT, "Give AI the paperwork", [
    [("Drafting: ", False, MUTED), ("3 min → 15 s", True, WHITE)],
    [("234 tickets ≈ ", False, MUTED), ("11 h saved", True, WHITE)],
    [("But the real win:", True, ACCENT)],
    [("Spring, by hand: ", False, MUTED), ("0% points", True, WHITE)],
    [("Agent-drafted: ", False, MUTED), ("90% points,", True, WHITE)],
    [("94% in the right epic", True, WHITE)],
], head_color=ACCENT)
footer(s, "forecasting needs points on every ticket")

NOTE_REFLECTION = (
    "Lesson 1 — we tried 3-day ticks from the agentic-augmented-scrum doc and moved to "
    "7-day cycles on Jira.\n\n"
    "Three reasons, and the third is the real one:\n"
    "  - too much changed inside a single tick to close it cleanly\n"
    "  - one small blocker consumed the whole cycle, because there was no slack\n"
    "  - reviewing an agent PR took longer than the agent took to write it\n\n"
    "The general point: agents moved the constraint from writing code to reviewing it. "
    "A 3-day cycle was sized for the old constraint. That is also why we plan capacity "
    "in review hours rather than story points.\n\n"
    "Lesson 2 — the AI-in-SE one. We gave agents control of ticket writing, "
    "documentation and comments, because they read the same repo we do and therefore "
    "carry more context than a person typing a ticket at the end of the day.\n\n"
    "Numbers, from dashboard/data/jira_issues.json (290 issues; the JQL and fetch time "
    "are recorded in the file):\n"
    "  - 3 minutes per ticket by hand, about 15 seconds for an agent\n"
    "  - 234 tickets created since 1 May, so roughly 11 hours saved\n"
    "  - that is under an hour a week. Do not oversell it.\n\n"
    "The number that matters is field completeness. Of the 56 tickets we wrote by hand "
    "in spring, ZERO had story points and ZERO had an epic parent. Of the 234 "
    "agent-drafted tickets, 90 percent have points and 94 percent have the right epic. "
    "A person under time pressure skips those fields; an agent does not.\n\n"
    "That is what makes the forecasting in the management section possible. Monte Carlo "
    "over closed-issue throughput needs points on every ticket. In spring we could not "
    "have produced that chart from our own backlog.\n\n"
    "Honest caveat if pushed: we review every agent-drafted ticket, and the 10 percent "
    "without points are mostly ones we corrected or closed as duplicates."
)
prs.slides[0].notes_slide.notes_text_frame.text = NOTE_REFLECTION
save(prs, OUT_REFLECTION, 1)
